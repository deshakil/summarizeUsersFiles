

from flask import Flask, request, jsonify, Response, stream_with_context
from flask_cors import CORS
from openai import AzureOpenAI
import os
import tempfile
import logging
import mammoth
import PyPDF2
import pandas as pd
from pptx import Presentation

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)

# Temporary storage for extracted text per session
document_cache = {}

def extract_text_from_file(file_path, file_type):
    """Extracts text from different file formats."""
    try:
        logger.info(f"Extracting text from {file_type} file: {file_path}")
        
        if file_type == "pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = "\n".join([page.extract_text() or "" for page in reader.pages])
            return text.strip()

        elif file_type == "docx":
            with open(file_path, "rb") as f:
                raw_text = mammoth.extract_raw_text(f)
            return raw_text.value.strip()

        elif file_type == "pptx":
            presentation = Presentation(file_path)
            text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
            return text.strip()

        elif file_type == "xlsx":
            df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
            text = "\n".join([df[sheet].to_csv(index=False) for sheet in df])
            return text.strip()

        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    except Exception as e:
        logger.error(f"Error extracting text from {file_type}: {str(e)}")
        raise ValueError(f"Failed to extract text from {file_type}: {str(e)}")


@app.route('/summarize', methods=['POST'])
def summarize():
    """Receives a file upload, extracts text, and sends it to OpenAI API."""
    try:
        logger.info("Received summarize request")
        
        # Validate request
        if 'file' not in request.files:
            logger.warning("No file uploaded")
            return jsonify({"error": "No file uploaded"}), 400

        uploaded_file = request.files['file']
        file_name = uploaded_file.filename
        
        if not file_name:
            logger.warning("Empty filename")
            return jsonify({"error": "Invalid filename"}), 400
            
        file_type = file_name.split('.')[-1].lower()
        logger.info(f"Processing file: {file_name}, type: {file_type}")

        # Validate file type
        allowed_types = {"pdf", "docx", "pptx", "xlsx"}
        if file_type not in allowed_types:
            logger.warning(f"Unsupported file type: {file_type}")
            return jsonify({"error": f"Unsupported file type: {file_type}"}), 400

        # Save file to a temporary location
        temp_file_path = os.path.join(tempfile.gettempdir(), file_name)
        uploaded_file.save(temp_file_path)
        logger.info(f"File saved to temporary path: {temp_file_path}")

        # Extract text from file
        extracted_text = extract_text_from_file(temp_file_path, file_type)
        logger.info(f"Extracted {len(extracted_text)} characters of text")

        # Ensure valid text extraction
        if not extracted_text.strip():
            logger.warning("No text extracted from file")
            return jsonify({"error": "No text extracted from file"}), 400

        # Store extracted text for follow-up questions
        session_id = request.headers.get('X-Session-ID', 'default')
        document_cache[session_id] = extracted_text
        
        # Validate API credentials
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            logger.error("Missing OpenAI API key")
            return jsonify({"error": "Server configuration error: Missing API key"}), 500
        
        # Create streaming response
        def generate():
            try:
                # Create client with API key
                client = AzureOpenAI(
                    api_key=api_key,
                    api_version="2024-12-01-preview",
                    azure_endpoint="https://weez-openai-resource.openai.azure.com/"
                )
            
                # Deployment Name (from Azure)
                DEPLOYMENT_NAME = "gpt-4o"  # Change to "gpt-4o" if needed
                logger.info(f"Sending request to OpenAI with model {DEPLOYMENT_NAME}")
                
                # Use streaming for the completion
                stream = client.chat.completions.create(
                    model=DEPLOYMENT_NAME,
                    messages=[
                        {"role": "system", "content": "Summarize this document in 3-5 bullet points:"},
                        {"role": "user", "content": extracted_text[:15000]}  # Limit input size
                    ],
                    temperature=0.3,
                    stream=True
                )
                
                for chunk in stream:
                    if chunk.choices and chunk.choices[0].delta.content:
                        content = chunk.choices[0].delta.content
                        yield content
                
                # Clean up temp file
                try:
                    os.remove(temp_file_path)
                except Exception as e:
                    logger.warning(f"Failed to remove temp file: {e}")
                    
            except Exception as e:
                logger.error(f"Error during streaming: {str(e)}", exc_info=True)
                yield f"Error generating summary: {str(e)}"
        
        # Return a streaming response
        return Response(stream_with_context(generate()), content_type='text/plain')
    
    except Exception as e:
        logger.error(f"Error in summarize endpoint: {str(e)}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route('/ask', methods=['POST'])
def ask():
    """Handles follow-up questions about the previously uploaded document."""
    try:
        logger.info("Received ask request")
        
        session_id = request.headers.get('X-Session-ID', 'default')
        
        # Ensure there's a stored document
        if session_id not in document_cache:
            logger.warning(f"No document found for session {session_id}")
            return jsonify({"error": "No document found. Please upload a file first."}), 400

        # Get the user's query
        data = request.get_json()
        if not data:
            logger.warning("No JSON data in request")
            return jsonify({"error": "Missing request data"}), 400
            
        query = data.get("query")

        if not query or not query.strip():
            logger.warning("Empty query received")
            return jsonify({"error": "Query cannot be empty"}), 400

        # Retrieve the stored document text
        document_text = document_cache[session_id]
        logger.info(f"Retrieved document ({len(document_text)} chars) for query: {query[:50]}...")

        # Validate API credentials
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            logger.error("Missing OpenAI API key")
            return jsonify({"error": "Server configuration error: Missing API key"}), 500

        # Create streaming response
        def generate():
            try:
                # Send follow-up question to OpenAI with the document context
                client = AzureOpenAI(
                    api_key=api_key,
                    api_version="2024-12-01-preview",
                    azure_endpoint="https://weez-openai-resource.openai.azure.com/"
                )
            
                # Deployment Name (from Azure)
                DEPLOYMENT_NAME = "gpt-4o"  # Change to "gpt-4o" if needed
                logger.info(f"Sending request to OpenAI with model {DEPLOYMENT_NAME}")
                
                # Use streaming for the completion
                stream = client.chat.completions.create(
                    model=DEPLOYMENT_NAME,
                    messages=[
                        {"role": "system", "content": "You are answering questions based on the following document:"},
                        {"role": "user", "content": document_text[:15000]},  # Limit input size
                        {"role": "user", "content": f"Based on this document, {query}"}
                    ],
                    temperature=0.3,
                    stream=True
                )
                
                for chunk in stream:
                    if chunk.choices and chunk.choices[0].delta.content:
                        content = chunk.choices[0].delta.content
                        yield content
                        
            except Exception as e:
                logger.error(f"Error during streaming: {str(e)}", exc_info=True)
                yield f"Error answering question: {str(e)}"
        
        # Return a streaming response
        return Response(stream_with_context(generate()), content_type='text/plain')

    except Exception as e:
        logger.error(f"Error in ask endpoint: {str(e)}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route('/health', methods=['GET'])
def health_check():
    """Simple endpoint to verify the API is running."""
    return jsonify({"status": "ok"})


if __name__ == '__main__':
    logger.info("Starting Flask server")
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port, debug=False)  # Disable debug in production
